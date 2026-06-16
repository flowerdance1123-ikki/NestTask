from pathlib import Path
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt
from pptx.dml.color import RGBColor

from .qr_generator import make_qr_image
from .image_utils import prepare_image_for_ppt, calc_fit_rect


LEFT_ALIGN_KEYS = {
    "body_status", "recent_state", "notes", "shared_points",
    "photo_1_caption", "photo_2_caption",
    "main_title", "main_url", "sub_title_1", "sub_url_1",
}

BLACK_TEXT_KEYS = {"patient_name", "therapist_name", "created_date"}

FALLBACK_FONT_SIZE_PT = {
    "body_status": 14, "recent_state": 14,
    "notes": 14, "shared_points": 14,
    "photo_1_caption": 14, "photo_2_caption": 14,
    "patient_name": 16, "therapist_name": 16, "created_date": 16,
    "main_title": 14, "main_url": 10,
    "sub_title_1": 14, "sub_url_1": 10,
}

QR_TAG_URL_KEYS = {
    "{{main_qr_box}}": ["main_qr_url", "main_url"],
    "{{sub_qr_box}}": ["sub_qr_url", "sub_url_1"],
}

PHOTO_TAGS = {"{{photo_1_box}}", "{{photo_2_box}}"}
PHOTO_TAG_PATH_KEYS = {
    "{{photo_1_box}}": "photo_1_path",
    "{{photo_2_box}}": "photo_2_path",
}

CLEAR_ONLY_TAGS = {"{{thumbnail_box}}"}

_TAG_RE = re.compile(r"\{\{[^}]+\}\}")


def find_remaining_tags(prs) -> list:
    """保存前の Presentation 内に残っている {{...}} タグを返す。"""
    found = set()
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                found.update(_TAG_RE.findall(shape.text_frame.text))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        found.update(_TAG_RE.findall(cell.text_frame.text))
    return sorted(found)


def normalize_patient_name(value: str) -> str:
    text = str(value).strip()
    text = re.sub(r"[ 　]+", "　", text)
    return text


def normalize_replacements(replacements: dict) -> dict:
    replacements = dict(replacements)
    if "patient_name" in replacements:
        replacements["patient_name"] = normalize_patient_name(replacements["patient_name"])
    return replacements


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def detect_replaced_keys(text: str, replacements: dict) -> set:
    if not text:
        return set()
    return {k for k in replacements if "{{" + k + "}}" in text}


def replace_text_in_string(text: str, replacements: dict):
    if not text:
        return text, False
    new_text = text
    for key, value in replacements.items():
        new_text = new_text.replace("{{" + key + "}}", str(value))
    return new_text, new_text != text


def get_first_font_size(text_frame):
    for para in text_frame.paragraphs:
        for run in para.runs:
            if run.text and run.font.size is not None:
                return run.font.size
    return None


def get_first_alignment(text_frame):
    for para in text_frame.paragraphs:
        if para.alignment is not None:
            return para.alignment
    return None


def decide_font_size(replaced_keys: set, captured_size):
    if captured_size is not None:
        return captured_size
    for key in replaced_keys:
        if key in FALLBACK_FONT_SIZE_PT:
            return Pt(FALLBACK_FONT_SIZE_PT[key])
    return None


def apply_font_size(text_frame, font_size) -> None:
    if font_size is None:
        return
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.size = font_size


def apply_font_color(text_frame, replaced_keys: set) -> None:
    if not (replaced_keys & BLACK_TEXT_KEYS):
        return
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0, 0, 0)


def apply_textbox_layout(text_frame, replaced_keys: set, original_alignment) -> None:
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE

    if replaced_keys & LEFT_ALIGN_KEYS:
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        for para in text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
    elif original_alignment is not None:
        for para in text_frame.paragraphs:
            para.alignment = original_alignment


def replace_tags_in_text_frame(text_frame, replacements: dict) -> int:
    original_text = text_frame.text
    replaced_keys = detect_replaced_keys(original_text, replacements)
    if not replaced_keys:
        return 0

    captured_font_size = get_first_font_size(text_frame)
    original_alignment = get_first_alignment(text_frame)
    new_text, changed = replace_text_in_string(original_text, replacements)

    if changed:
        text_frame.text = new_text
        font_size = decide_font_size(replaced_keys, captured_font_size)
        apply_font_size(text_frame, font_size)
        apply_font_color(text_frame, replaced_keys)
        apply_textbox_layout(text_frame, replaced_keys, original_alignment)
        return 1
    return 0


def replace_tags_in_shape(shape, replacements: dict) -> int:
    count = 0
    if getattr(shape, "has_text_frame", False):
        count += replace_tags_in_text_frame(shape.text_frame, replacements)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                count += replace_tags_in_text_frame(cell.text_frame, replacements)
    return count


def _get_first_available_url(replacements: dict, url_keys: list) -> str:
    for key in url_keys:
        value = str(replacements.get(key, "")).strip()
        if value:
            return value
    return ""


def insert_media_images(prs, replacements: dict, temp_dir: Path) -> dict:
    temp_qr_dir = temp_dir / "qr"
    temp_photo_dir = temp_dir / "photo"
    stats = {"qr": 0, "photo": 0, "cleared": 0}

    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue

            text = shape.text.strip()

            if text in QR_TAG_URL_KEYS:
                url = _get_first_available_url(replacements, QR_TAG_URL_KEYS[text])
                shape.text_frame.text = ""

                if url:
                    tag_stem = text.strip("{}").replace("_box", "")
                    qr_path = temp_qr_dir / f"{tag_stem}.png"
                    make_qr_image(url, qr_path)
                    slide.shapes.add_picture(
                        str(qr_path),
                        shape.left, shape.top,
                        width=shape.width, height=shape.height,
                    )
                    stats["qr"] += 1
                else:
                    stats["cleared"] += 1

            elif text in PHOTO_TAGS:
                path_key = PHOTO_TAG_PATH_KEYS[text]
                photo_path_str = str(replacements.get(path_key, "")).strip()
                photo_path = Path(photo_path_str) if photo_path_str else None
                shape.text_frame.text = ""

                if photo_path and photo_path.exists():
                    stem = text.strip("{}").replace("_box", "")
                    prepared_path, img_w, img_h = prepare_image_for_ppt(
                        photo_path, temp_photo_dir, stem
                    )
                    left, top, width, height = calc_fit_rect(
                        shape.left, shape.top,
                        shape.width, shape.height,
                        img_w, img_h,
                    )
                    slide.shapes.add_picture(
                        str(prepared_path), left, top, width=width, height=height
                    )
                    stats["photo"] += 1
                else:
                    stats["cleared"] += 1

            elif text in CLEAR_ONLY_TAGS:
                shape.text_frame.text = ""
                stats["cleared"] += 1

    return stats


def generate_pptx(
    template_path: Path,
    replacements: dict,
    output_path: Path,
    temp_dir: Path,
) -> dict:
    """テンプレートにデータを差し込みPPTXを生成する。辞書で統計情報を返す。"""
    if not template_path.exists():
        raise FileNotFoundError(
            f"テンプレートファイルが見つかりません。\n{template_path.name} があるか確認してください。"
        )

    replacements = normalize_replacements(replacements)
    prs = Presentation(str(template_path))

    replaced_count = 0
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            replaced_count += replace_tags_in_shape(shape, replacements)

    stats = insert_media_images(prs, replacements, temp_dir)
    stats["text_replaced"] = replaced_count

    output_path.parent.mkdir(parents=True, exist_ok=True)
    stats["unresolved_tags"] = find_remaining_tags(prs)
    prs.save(str(output_path))

    return stats
